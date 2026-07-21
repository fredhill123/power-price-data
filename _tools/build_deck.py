"""
build_deck.py — build HourlyPowerData.pptx: a Rothschild-style deck (on the P&U
Crash Course template) whose 9 charts are LINKED to HourlyPowerData.xlsx, so they
auto-update on refresh.

Two phases:
  1) python-pptx: open the crash-course template (inherit master/layouts/theme/
     Calluna+Arial fonts + R&Co colours), strip its slides, add a title slide +
     6 content slides with title/kicker/navy exhibit bars/captions/source line.
  2) zip surgery: inject the 9 Redburn chart parts as LINKED charts
     (<c:externalData> -> oleObject external rel -> the UNC path of the .xlsx),
     plus the <p:graphicFrame> on each slide. No embedded workbook — a pure link,
     exactly like the team's existing linked decks.

Usage: python build_deck.py <template.pptx> <charts_source.xlsx> <out.pptx>
"""
from __future__ import annotations
import sys, os, re, zipfile, shutil, tempfile
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- house colours / fonts ----
NAVY_BAR = RGBColor(0x1C, 0x35, 0x5E)   # exhibit title bar (dk2)
BRAND    = RGBColor(0x2E, 0x3E, 0x80)   # title / kicker (accent1)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
SRC_GREY = RGBColor(0x9A, 0x9A, 0x9A)
SERIF = "Calluna"; SANS = "Arial"

# ---- external link (UNC form; H:\ maps to \\redburn.local\core\data\) ----
LINK_TARGET = (r"file:///\\redburn.local\core\data\Oils\Oils%202.0"
               r"\Power%20%26%20Utilities%20Team%20Resources"
               r"\Sector%20Presentation\HourlyPowerData.xlsx")

C="http://schemas.openxmlformats.org/drawingml/2006/chart"
def cc(t): return f"{{{C}}}{t}"
ORDER_SPACE=["date1904","lang","roundedCorners","style","clrMapOvr","pivotSource",
             "protection","chart","spPr","txPr","externalData","printSettings",
             "userShapes","extLst"]
def insert_ordered(parent, child, order):
    tag=etree.QName(child).localname
    idx=order.index(tag) if tag in order else len(order)
    pos=len(parent)
    for i,ex in enumerate(parent):
        ln=etree.QName(ex).localname
        if ln in order and order.index(ln) > idx: pos=i; break
    parent.insert(pos, child)

# ---- content geometry (EMU) ----
LEFT=596900; FULLW=8712200
BAR_TOP=1350000; BAR_H=260000
CH_TOP=1660000
CH_H_2UP=4450000; CH_H_1UP=4550000
COL_L=596900; COL_R=5059000; COLW=4249100   # 2-up columns

# ---- the deck plan ----
# each content slide: (title, kicker, [ (chart_src_num, caption, (x,y,w,h)) ... ])
def box_1up():  return (LEFT, CH_TOP, FULLW, CH_H_1UP)
def bar_1up():  return (LEFT, BAR_TOP, FULLW, BAR_H)
def box_L():    return (COL_L, CH_TOP, COLW, CH_H_2UP)
def box_R():    return (COL_R, CH_TOP, COLW, CH_H_2UP)
def bar_L():    return (COL_L, BAR_TOP, COLW, BAR_H)
def bar_R():    return (COL_R, BAR_TOP, COLW, BAR_H)

SOURCE = "Source: ENTSO-E Transparency Platform (hourly prices, generation, flows & capacity)."

# ---- slides derived from deck_spec.py (THE single source of truth) ----
# Linked-chart slides = every exhibit has a workbook chart#; snapshot slides = image-only.
import deck_spec
_BOXFN = {"L": box_L, "R": box_R, "1up": box_1up}
_BARFN = {"L": bar_L, "R": bar_R, "1up": bar_1up}

def _linked_slides():
    out = []
    for s in deck_spec.SLIDES:
        if all(e["chart"] is not None for e in s["exhibits"]):
            out.append((s["title"], s["kicker"],
                        [(e["chart"], e["caption"], _BOXFN[e["box"]](), _BARFN[e["box"]]())
                         for e in s["exhibits"]]))
    return out

def _static_slides():
    out = []
    for s in deck_spec.SLIDES:
        if all(e["chart"] is None for e in s["exhibits"]):
            out.append((s["title"], s["kicker"],
                        [(e["png"], e["caption"], _BOXFN[e["box"]](), _BARFN[e["box"]]())
                         for e in s["exhibits"]]))
    return out

_OUTPUTS = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "outputs")
IMG_DIR = os.path.join(_OUTPUTS, "deck_img")        # pre-rendered snapshots (snap_*)
CHARTS_DIR = os.path.join(_OUTPUTS, "deck_charts")  # render_all output (everything else, e.g. F scatter)
SLIDES = _linked_slides()
STATIC_SLIDES = _static_slides()

def set_ph_text(slide, idx, text, font=SERIF, size=None, color=None, bold=False):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx==idx:
            tf=ph.text_frame; tf.text=text
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.name=font
                    if size: r.font.size=Pt(size)
                    if color is not None: r.font.color.rgb=color
                    r.font.bold=bold
            return ph
    return None

def add_bar(slide, box, caption):
    x,y,w,h=box
    rect=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x),Emu(y),Emu(w),Emu(h))
    rect.fill.solid(); rect.fill.fore_color.rgb=NAVY_BAR
    rect.line.fill.background()
    rect.shadow.inherit=False
    tf=rect.text_frame; tf.word_wrap=True
    tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Emu(64000); tf.margin_right=Emu(45000)
    tf.margin_top=Emu(0); tf.margin_bottom=Emu(0)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    run=p.add_run(); run.text=caption
    run.font.name=SANS; run.font.size=Pt(8.5); run.font.bold=True; run.font.color.rgb=WHITE
    return rect

def build_base(template, out_base):
    prs=Presentation(template)
    # strip all slides (drop the relationship so the part isn't written)
    sldIdLst=prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId=sldId.get(qn("r:id")); sldIdLst.remove(sldId); prs.part.drop_rel(rId)

    # ---- title slide (layout 0) ----
    title=prs.slides.add_slide(prs.slide_layouts[0])
    set_ph_text(title, 0, "Hourly Power Data", SERIF, 40, BRAND)
    set_ph_text(title, 1, "European wholesale power — prices, generation & capacity", SERIF, 20, BRAND)
    set_ph_text(title, 13, "17 July 2026", SANS, 12)
    set_ph_text(title, 14, "ENTSO-E hourly data · auto-linked workbook", SANS, 9, SRC_GREY)

    # ---- content slides (layout 12 'Title Only') ----
    for stitle, kicker, charts in SLIDES:
        s=prs.slides.add_slide(prs.slide_layouts[12])
        set_ph_text(s, 0, stitle, SERIF, 22.5, BRAND)
        set_ph_text(s, 13, kicker, SERIF, 13.1, BRAND)
        # source line: reposition to bottom-LEFT (matches crash course; clears the
        # template's centred "Confidential" stamp, which the 'Title Only' layout
        # would otherwise overlap by placing source on the right).
        src=set_ph_text(s, 16, SOURCE, SANS, 7, SRC_GREY)
        if src is not None:
            src.left=Emu(596900); src.top=Emu(6640000)
            src.width=Emu(4300000); src.height=Emu(300000)
            for p in src.text_frame.paragraphs: p.alignment=PP_ALIGN.LEFT
        for (_num, caption, _box, bar) in charts:
            add_bar(s, bar, caption)

    # ---- static-image content slides ----
    for stitle, kicker, imgs in STATIC_SLIDES:
        s=prs.slides.add_slide(prs.slide_layouts[12])
        set_ph_text(s, 0, stitle, SERIF, 22.5, BRAND)
        set_ph_text(s, 13, kicker, SERIF, 13.1, BRAND)
        src=set_ph_text(s, 16, SOURCE, SANS, 7, SRC_GREY)
        if src is not None:
            src.left=Emu(596900); src.top=Emu(6640000)
            src.width=Emu(4300000); src.height=Emu(300000)
            for p in src.text_frame.paragraphs: p.alignment=PP_ALIGN.LEFT
        for (png, caption, box, bar) in imgs:
            add_bar(s, bar, caption)
            x,y,w,h=box
            # snapshots live in deck_img; render_all exhibits (e.g. F scatter) in deck_charts
            path=os.path.join(IMG_DIR if png.startswith("snap_") else CHARTS_DIR, png)
            # width-fit; height scales to preserve the image's aspect ratio
            s.shapes.add_picture(path, Emu(x), Emu(y), width=Emu(w))
    prs.save(out_base)

# ---------- phase 2: inject linked charts ----------
def add_externaldata(chart_xml_bytes, rid):
    root=etree.fromstring(chart_xml_bytes)
    for old in root.findall(cc("externalData")): root.remove(old)
    ed=etree.Element(cc("externalData")); ed.set(f"{{http://schemas.openxmlformats.org/officeDocument/2006/relationships}}id", rid)
    etree.SubElement(ed, cc("autoUpdate")).set("val","0")
    insert_ordered(root, ed, ORDER_SPACE)
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

def chart_rels():
    return ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            f'<Relationship Id="rId1" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject" '
            f'Target="{LINK_TARGET}" TargetMode="External"/></Relationships>')

def graphicframe_xml(shape_id, name, rid, box):
    x,y,w,h=box
    return (f'<p:graphicFrame><p:nvGraphicFramePr>'
            f'<p:cNvPr id="{shape_id}" name="{name}"/>'
            f'<p:cNvGraphicFramePr><a:graphicFrameLocks/></p:cNvGraphicFramePr><p:nvPr/>'
            f'</p:nvGraphicFramePr>'
            f'<p:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></p:xfrm>'
            f'<a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart">'
            f'<c:chart xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
            f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" r:id="{rid}"/>'
            f'</a:graphicData></a:graphic></p:graphicFrame>')

def inject(base_pptx, xlsx_src, out_pptx):
    # read all base parts
    zin=zipfile.ZipFile(base_pptx)
    parts={i.filename: zin.read(i.filename) for i in zin.infolist()}
    order=[i.filename for i in zin.infolist()]; zin.close()
    # source chart XMLs from the workbook
    zx=zipfile.ZipFile(xlsx_src)
    src_chart={n: zx.read(n) for n in zx.namelist() if re.match(r"xl/charts/chart\d+\.xml$", n)}
    zx.close()

    # content slides were added in order -> slide2..slide7 (slide1 = title)
    chart_no=0
    for si, (stitle, kicker, charts) in enumerate(SLIDES, start=2):
        slide_file=f"ppt/slides/slide{si}.xml"
        rel_file=f"ppt/slides/_rels/slide{si}.xml.rels"
        sx=parts[slide_file].decode()
        rels=parts[rel_file].decode()
        # next free slide-local rId
        used=[int(m) for m in re.findall(r'Id="rId(\d+)"', rels)]
        next_rid=(max(used)+1) if used else 1
        # next free shape id on the slide
        sid=max([int(m) for m in re.findall(r'<p:cNvPr id="(\d+)"', sx)] or [1])
        gframes=""
        for (num, caption, box, bar) in charts:
            chart_no+=1; sid+=1
            cpart=f"ppt/charts/chart{chart_no}.xml"
            cpart_rels=f"ppt/charts/_rels/chart{chart_no}.xml.rels"
            # chart part = redburn chart XML + externalData(rId1)
            parts[cpart]=add_externaldata(src_chart[f"xl/charts/chart{num}.xml"], "rId1")
            parts[cpart_rels]=chart_rels().encode()
            order += [cpart, cpart_rels]
            # slide rel -> chart part
            rid=f"rId{next_rid}"; next_rid+=1
            rels=rels.replace("</Relationships>",
                f'<Relationship Id="{rid}" '
                'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart" '
                f'Target="../charts/chart{chart_no}.xml"/></Relationships>')
            gframes+=graphicframe_xml(sid, f"Chart {chart_no}", rid, box)
        parts[slide_file]=sx.replace("</p:spTree>", gframes+"</p:spTree>").encode()
        parts[rel_file]=rels.encode()

    # content types: add chart overrides
    ct=parts["[Content_Types].xml"].decode()
    adds="".join(f'<Override PartName="/ppt/charts/chart{k}.xml" '
                 'ContentType="application/vnd.openxmlformats-officedocument.drawingml.chart+xml"/>'
                 for k in range(1, chart_no+1))
    parts["[Content_Types].xml"]=ct.replace("</Types>", adds+"</Types>").encode()

    zout=zipfile.ZipFile(out_pptx,"w",zipfile.ZIP_DEFLATED)
    for name in order: zout.writestr(name, parts[name])
    zout.close()
    print(f"injected {chart_no} linked charts -> {out_pptx}", flush=True)

def main():
    template, xlsx_src, out_pptx = sys.argv[1], sys.argv[2], sys.argv[3]
    base=tempfile.mktemp(suffix=".pptx")
    build_base(template, base)
    inject(base, xlsx_src, out_pptx)
    os.remove(base)

if __name__=="__main__":
    main()
