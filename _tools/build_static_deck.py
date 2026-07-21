"""
build_static_deck.py — assemble the Claude-rendered PNGs (render_all.py +
render_snapshots.py) into a self-contained Rothschild/Redburn deck. NO Excel, NO
linked charts — every exhibit is an embedded image, so the file is a frozen
snapshot of the day it was built.

Mirrors the linked deck (build_deck.py) slide-for-slide, caption-for-caption, so
the two stay visually identical — same title/kicker/navy exhibit bars/source line.

Usage: python build_static_deck.py <out.pptx>   (template + geometry from build_deck)
"""
from __future__ import annotations
import os, sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
import build_deck as bd          # reuse branding + geometry + helpers
import completeness
import deck_spec                  # SINGLE SOURCE OF TRUTH

CHARTS = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "outputs", "deck_charts")
IMG    = bd.IMG_DIR              # outputs/deck_img (Spain snapshots)
TEMPLATE = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                        "archive", "phase4_2026-07-17", "HourlyPowerData_pre-phase4.pptx")

_BOX = {"L": bd.box_L, "R": bd.box_R, "1up": bd.box_1up}
_BAR = {"L": bd.bar_L, "R": bd.bar_R, "1up": bd.bar_1up}

def png_path(name):
    return os.path.join(IMG if name.startswith("snap_") else CHARTS, name)


def build(out_pptx):
    ce = completeness.cutoffs()["coverage_end"]
    as_of = ce.strftime("%-d %B %Y")
    prs = Presentation(TEMPLATE)
    # strip template slides
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        rid = sid.get(bd.qn("r:id")); lst.remove(sid); prs.part.drop_rel(rid)
    # title
    t = prs.slides.add_slide(prs.slide_layouts[0])
    bd.set_ph_text(t, 0, "Hourly Power Data", bd.SERIF, 40, bd.BRAND)
    bd.set_ph_text(t, 1, "European wholesale power — prices, generation & capacity", bd.SERIF, 20, bd.BRAND)
    bd.set_ph_text(t, 13, f"Data as of {as_of}", bd.SANS, 12)
    bd.set_ph_text(t, 14, "ENTSO-E hourly data · self-contained snapshot (no live links)", bd.SANS, 9, bd.SRC_GREY)
    # content (driven by deck_spec — the single source of truth)
    for slide in deck_spec.SLIDES:
        s = prs.slides.add_slide(prs.slide_layouts[12])
        bd.set_ph_text(s, 0, slide["title"], bd.SERIF, 22.5, bd.BRAND)
        bd.set_ph_text(s, 13, slide["kicker"], bd.SERIF, 13.1, bd.BRAND)
        src = bd.set_ph_text(s, 16, bd.SOURCE, bd.SANS, 7, bd.SRC_GREY)
        if src is not None:
            src.left = Emu(596900); src.top = Emu(6640000)
            src.width = Emu(4300000); src.height = Emu(300000)
            for p in src.text_frame.paragraphs: p.alignment = PP_ALIGN.LEFT
        for e in slide["exhibits"]:
            bd.add_bar(s, _BAR[e["box"]](), e["caption"])
            x, y, w, h = _BOX[e["box"]]()
            s.shapes.add_picture(png_path(e["png"]), Emu(x), Emu(y), width=Emu(w))
    prs.save(out_pptx)
    print(f"built static deck ({len(deck_spec.SLIDES)+1} slides, data as of {as_of}) -> {out_pptx}", flush=True)


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "outputs", "HourlyPowerData_snapshot.pptx")
    build(out)
