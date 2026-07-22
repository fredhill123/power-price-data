"""
check_consistency.py — fail-loud guard that the LINKED deck and the STATIC deck
never drift from deck_spec.py (or each other). Run before every delivery.

Asserts, for BOTH decks:
  * slide count == 1 (title) + len(deck_spec.SLIDES)
  * each content slide's title + kicker == deck_spec
  * each content slide's navy-bar captions (in order) == deck_spec exhibit captions
And that the linked workbook holds charts 1..19.

Exit 0 = consistent; exit 1 = drift (prints the diffs).
"""
from __future__ import annotations
import os, sys, zipfile, re
from xml.etree import ElementTree as ET
from pptx import Presentation
import deck_spec

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LINKED = os.path.join(ROOT, "outputs", "HourlyPowerData.pptx")
STATIC = os.path.join(ROOT, "outputs", "HourlyPowerData_snapshot.pptx")
WORKBOOK = os.path.join(ROOT, "outputs", "HourlyPowerData.xlsx")


def deck_content(path):
    """Return [(title, kicker, [captions...]), ...] for content slides (skip title)."""
    prs = Presentation(path)
    out = []
    for s in list(prs.slides)[1:]:
        title = kicker = ""
        caps = []
        for ph in s.placeholders:
            idx = ph.placeholder_format.idx
            if idx == 0: title = ph.text.strip()
            elif idx == 13: kicker = ph.text.strip()
        for sh in s.shapes:
            if not sh.is_placeholder and sh.has_text_frame and sh.text_frame.text.strip():
                caps.append(sh.text_frame.text.strip())
        out.append((title, kicker, caps))
    return out


def expected():
    return [(s["title"], s["kicker"], [e["caption"] for e in s["exhibits"]]) for s in deck_spec.SLIDES]


def check_deck(name, path, exp):
    errs = []
    if not os.path.exists(path):
        return [f"{name}: file missing ({path})"]
    got = deck_content(path)
    if len(got) != len(exp):
        errs.append(f"{name}: {len(got)} content slides, expected {len(exp)}")
    for i, (e, g) in enumerate(zip(exp, got), start=1):
        if e[0] != g[0]:
            errs.append(f"{name} slide {i} title: got {g[0]!r} != spec {e[0]!r}")
        if e[1] != g[1]:
            errs.append(f"{name} slide {i} kicker: got {g[1]!r} != spec {e[1]!r}")
        if e[2] != g[2]:
            errs.append(f"{name} slide {i} captions: got {g[2]} != spec {e[2]}")
    return errs


def check_xml_wellformed(path, label):
    """Every XML part must parse. This is the cheapest possible check and the one that
    matters most: a single unescaped '&' in one cell makes Excel offer to Recover the
    workbook, and Recovery strips Power Query. On 2026-07-22 a text column written into a
    numeric <v> shipped three malformed sheets through a PASSing consistency run, because
    nothing here had ever actually parsed the file it was signing off.
    """
    errs = []
    z = zipfile.ZipFile(path)
    for n in z.namelist():
        if not n.endswith((".xml", ".rels")):
            continue
        try:
            ET.fromstring(z.read(n))
        except ET.ParseError as e:
            errs.append(f"{label} {n}: malformed XML — {e}")
    return errs


def check_refresh_stability(path):
    """The two halves of the invariant that keeps charts intact through a refresh.

    Excel re-anchors a chart series whose range runs to or past the end of the data it
    reads. So (a) every query tab's pre-filled cache must be exactly the size of the CSV
    it will load, and (b) no chart range may extend past that cache. Break either and a
    chart silently re-fits itself on the user's first refresh — which is how chart12
    quietly re-acquired the six technologies curate_tech_charts.py had removed.
    """
    errs = []
    z = zipfile.ZipFile(path)
    parts = {n: z.read(n) for n in z.namelist()}
    wb = parts["xl/workbook.xml"].decode()
    relmap = dict(re.findall(r'Id="([^"]+)"[^>]*Target="([^"]+)"',
                             parts["xl/_rels/workbook.xml.rels"].decode()))
    extent = {}
    for name, rid in re.findall(r'<sheet name="([^"]+)"[^>]*?r:id="(rId\d+)"', wb):
        spart = "xl/" + relmap[rid].lstrip("/")
        sx = parts[spart].decode(errors="replace")
        d = re.search(r'<dimension ref="[A-Z]+\d+:[A-Z]+(\d+)"', sx)
        if d:
            extent[name] = int(d.group(1))
        rels = parts.get(spart.replace("worksheets/", "worksheets/_rels/") + ".rels", b"").decode()
        t = re.search(r'Target="\.\./tables/(table\d+\.xml)"', rels)
        if not t or name not in extent:
            continue
        tbl = parts["xl/tables/" + t.group(1)].decode()
        stem = re.search(r'\sname="([^"]+)"', tbl).group(1)
        csv_path = os.path.join(ROOT, "published", "charts", stem + ".csv")
        if not os.path.exists(csv_path):
            continue
        with open(csv_path, newline="") as f:
            want = sum(1 for _ in f)
        if extent[name] != want:
            errs.append(f"WORKBOOK {name}: pre-filled {extent[name]} rows but "
                        f"{stem}.csv has {want} — the table will change shape on refresh "
                        f"and Excel may re-anchor a chart (run resync_prefill.py)")

    for n in sorted(x for x in parts if re.match(r"xl/charts/chart\d+\.xml$", x)):
        cx = parts[n].decode()
        for f in sorted(set(re.findall(r"<c:f>([^<]+)</c:f>", cx))):
            m = re.match(r"([^!]+)!\$[A-Z]+\$\d+:\$[A-Z]+\$(\d+)$", f)
            if m and m.group(1) in extent and int(m.group(2)) > extent[m.group(1)]:
                errs.append(f"WORKBOOK {os.path.basename(n)}: range {f} runs past the "
                            f"pre-filled data (row {extent[m.group(1)]}) — Excel will "
                            f"stretch this series on refresh")
    return errs


def main():
    exp = expected()
    errs = []
    errs += check_deck("LINKED", LINKED, exp)
    errs += check_deck("STATIC", STATIC, exp)
    # workbook holds charts 1..15
    if os.path.exists(WORKBOOK):
        z = zipfile.ZipFile(WORKBOOK)
        nums = sorted(int(re.search(r"chart(\d+)", n).group(1))
                      for n in z.namelist() if re.match(r"xl/charts/chart\d+\.xml$", n))
        if nums != list(range(1, 20)):
            errs.append(f"WORKBOOK charts: {nums} != 1..19")
        errs += check_xml_wellformed(WORKBOOK, 'WORKBOOK')
        errs += check_refresh_stability(WORKBOOK)
    else:
        errs.append(f"WORKBOOK missing ({WORKBOOK})")
    frozen = os.path.join(ROOT, 'outputs', 'HourlyPowerData_frozen.xlsx')
    if os.path.exists(frozen):
        errs += check_xml_wellformed(frozen, 'FROZEN')

    if errs:
        print("CONSISTENCY: FAIL")
        for e in errs: print("  ✗", e)
        sys.exit(1)
    print(f"CONSISTENCY: PASS — both decks match deck_spec ({len(exp)} content slides, "
          f"{sum(len(s['exhibits']) for s in deck_spec.SLIDES)} exhibits) + workbook charts 1-19.")


if __name__ == "__main__":
    main()
